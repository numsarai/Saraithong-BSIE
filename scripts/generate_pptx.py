"""
generate_pptx.py — SPNI/BSIE Comprehensive Presentation
Proportional layout, consistent spacing, 22 slides.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Colors ──────────────────────────────────────────────────
NAVY    = RGBColor(0x1A, 0x36, 0x5D)
BLUE    = RGBColor(0x2B, 0x6C, 0xB0)
SKY     = RGBColor(0x63, 0xB3, 0xED)
GREEN   = RGBColor(0x38, 0xA1, 0x69)
DGREEN  = RGBColor(0x22, 0x54, 0x3D)
MINT    = RGBColor(0x48, 0xBB, 0x78)
RED     = RGBColor(0xE5, 0x3E, 0x3E)
DRED    = RGBColor(0xC5, 0x30, 0x30)
SALMON  = RGBColor(0xFC, 0x81, 0x81)
ORANGE  = RGBColor(0xDD, 0x6B, 0x20)
GOLD    = RGBColor(0xD6, 0x9E, 0x2E)
LGOLD   = RGBColor(0xEC, 0xC9, 0x4B)
PURPLE  = RGBColor(0x80, 0x5A, 0xD5)
LILAC   = RGBColor(0xD6, 0xBC, 0xFA)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY   = RGBColor(0xF7, 0xFA, 0xFC)
MGRAY   = RGBColor(0xED, 0xF2, 0xF7)
DGRAY   = RGBColor(0x2D, 0x37, 0x48)
VDARK   = RGBColor(0x1A, 0x20, 0x2C)
TXT     = RGBColor(0x1A, 0x20, 0x2C)
TXTL    = RGBColor(0x4A, 0x55, 0x68)
SUBTLE  = RGBColor(0xA0, 0xAE, 0xC0)
CDARK   = RGBColor(0x1E, 0x3A, 0x5F)

# ── Layout constants (Widescreen 16:9) ──────────────────────
SW = Inches(13.333)
SH = Inches(7.5)
MX = Inches(0.6)          # horizontal margin
CW = SW - MX * 2          # content width ~12.13"
GAP = Inches(0.25)        # gap between cards
TITLE_Y  = Inches(0.25)
TITLE_H  = Inches(0.65)
BODY_TOP = Inches(1.05)   # content starts below title
BODY_H   = SH - BODY_TOP - Inches(0.3)  # available height

def _half():
    """Width of a half-column card."""
    return (CW - GAP) / 2

def _third():
    """Width of a third-column card."""
    return (CW - GAP * 2) / 3


# ── Drawing primitives ──────────────────────────────────────

def _bg(slide, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()

def _rect(slide, l, t, w, h, color, border=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    if border:
        s.line.color.rgb = border; s.line.width = Pt(1.5)
    else:
        s.line.fill.background()
    return s

def _line(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def _txt(slide, l, t, w, h, text, sz=14, clr=TXT, bold=False, al=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = al
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.color.rgb = clr; r.font.bold = bold
    r.font.name = "Noto Sans Thai"
    return tb

def _title(slide, text, clr=NAVY):
    _txt(slide, MX, TITLE_Y, CW, TITLE_H, text, 28, clr, True, PP_ALIGN.CENTER)

def _sub(slide, text, clr=TXTL, y=None):
    yy = y or (TITLE_Y + Inches(0.55))
    _txt(slide, MX, yy, CW, Inches(0.35), text, 13, clr, False, PP_ALIGN.CENTER)

def _bullets(slide, l, t, w, h, items, sz=13, clr=TXT, bc=GREEN, sp=Pt(4)):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = sp; p.space_before = Pt(0)
        br = p.add_run(); br.text = "● "; br.font.size = Pt(sz - 2)
        br.font.color.rgb = bc; br.font.name = "Noto Sans Thai"
        parts = item.split("**")
        for j, part in enumerate(parts):
            if not part: continue
            r = p.add_run(); r.text = part; r.font.size = Pt(sz)
            r.font.color.rgb = clr; r.font.bold = (j % 2 == 1)
            r.font.name = "Noto Sans Thai"

def _card_title(slide, l, t, w, text, clr=NAVY):
    _txt(slide, l + Inches(0.2), t + Inches(0.12), w - Inches(0.4), Inches(0.35),
         text, 17, clr, True)


# ══════════════════════════════════════════════════════════════
def generate():
    prs = Presentation()
    prs.slide_width = SW; prs.slide_height = SH
    BL = prs.slide_layouts[6]

    def S(color=LGRAY):
        s = prs.slides.add_slide(BL); _bg(s, color); return s

    PAD = Inches(0.2)    # card internal padding
    hw = _half()          # half-width card
    tw = _third()         # third-width card

    # ════════════════════════════════════════════════════════════
    # 1. TITLE
    # ════════════════════════════════════════════════════════════
    s = S(NAVY)
    _txt(s, MX, Inches(0.8), CW, Inches(0.35),
         "สำนักงานตำรวจแห่งชาติ", 12, SUBTLE, False, PP_ALIGN.CENTER)
    _txt(s, MX, Inches(1.5), CW, Inches(0.9),
         "SPNI Platform", 52, WHITE, True, PP_ALIGN.CENTER)
    _txt(s, MX, Inches(2.5), CW, Inches(0.55),
         "Smart Police & National Intelligence", 24, SKY, False, PP_ALIGN.CENTER)
    _line(s, Inches(5.9), Inches(3.3), Inches(1.5), Pt(3), GOLD)
    _txt(s, MX, Inches(3.6), CW, Inches(0.45),
         "โมดูลแรก: BSIE v4.0 — Bank Statement Intelligence Engine", 20, WHITE, False, PP_ALIGN.CENTER)
    _txt(s, MX, Inches(4.2), CW, Inches(0.4),
         "ระบบวิเคราะห์ธุรกรรมทางการเงินอัจฉริยะ สำหรับงานสืบสวนสอบสวน", 16, RGBColor(0xBE,0xE3,0xF8), False, PP_ALIGN.CENTER)
    _txt(s, MX, Inches(5.5), CW, Inches(1.2),
         "เอกสารประกอบการนำเสนอเพื่อของบประมาณระดับประเทศ\n"
         "จัดทำโดย ร้อยตำรวจเอกณัฐวุฒิ สาหร่ายทอง\nเมษายน 2569",
         13, SUBTLE, False, PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════════
    # 2. SPNI VISION
    # ════════════════════════════════════════════════════════════
    s = S(NAVY)
    _title(s, "ระบบ SPNI — ภาพรวม 4 ส่วนงาน", WHITE)
    _sub(s, "ระบบรวมเครื่องมือช่วยงานสืบสวนสอบสวน สำหรับ สตช. ทั้งประเทศ", RGBColor(0xBE,0xE3,0xF8))

    mods = [
        ("โมดูล 1", "วิเคราะห์การเงิน", "v4.0 พร้อมใช้ ★", MINT),
        ("โมดูล 2", "วิเคราะห์โทรศัพท์", "วางแผน", DGRAY),
        ("โมดูล 3", "สื่อสังคมออนไลน์", "วางแผน", DGRAY),
        ("โมดูล 4", "กล้องวงจรปิด", "วางแผน", DGRAY),
    ]
    mw = (CW - GAP * 3) / 4
    for i, (num, name, status, bg_c) in enumerate(mods):
        x = MX + i * (mw + GAP)
        bc = GOLD if i == 0 else RGBColor(0x4A,0x55,0x68)
        _rect(s, x, Inches(1.55), mw, Inches(1.6), bg_c, bc)
        _txt(s, x, Inches(1.6), mw, Inches(0.3), num, 10, SUBTLE, False, PP_ALIGN.CENTER)
        _txt(s, x, Inches(1.9), mw, Inches(0.4), name, 15, WHITE, True, PP_ALIGN.CENTER)
        _txt(s, x, Inches(2.35), mw, Inches(0.3), status, 11, GOLD if i==0 else SUBTLE, False, PP_ALIGN.CENTER)
        if i < 3:
            _txt(s, x + mw + Inches(0.02), Inches(2.0), Inches(0.2), Inches(0.3), "→", 16, SKY, False, PP_ALIGN.CENTER)

    # Shared layer
    _rect(s, MX, Inches(3.45), CW, Inches(1.6), CDARK)
    _txt(s, MX, Inches(3.5), CW, Inches(0.35), "ชั้นแบ่งปันข่าวกรองกลาง", 16, WHITE, True, PP_ALIGN.CENTER)
    tags = [("ทะเบียนบุคคล", GREEN), ("วิเคราะห์ความเชื่อมโยง", BLUE), ("จัดการคดี", ORANGE),
            ("AI ภายในเครื่อง", RED), ("สายโซ่หลักฐาน", PURPLE), ("รายงาน กม.", GOLD)]
    tw_tag = (CW - Inches(0.6) - GAP * 5) / 6
    for i, (name, clr) in enumerate(tags):
        x = MX + Inches(0.3) + i * (tw_tag + GAP)
        _rect(s, x, Inches(4.0), tw_tag, Inches(0.45), clr)
        _txt(s, x, Inches(4.02), tw_tag, Inches(0.42), name, 10, WHITE, True, PP_ALIGN.CENTER)

    _txt(s, MX, Inches(4.65), CW, Inches(0.35),
         "ทุกส่วนงานเชื่อมข้อมูลกัน — ค้นพบคนเดียวกันจากบัญชีธนาคาร + โทรศัพท์ + สื่อสังคม + กล้องวงจรปิด", 12, RGBColor(0xBE,0xE3,0xF8), False, PP_ALIGN.CENTER)

    # Infra
    _rect(s, MX, Inches(5.3), CW, Inches(0.55), DGRAY)
    _txt(s, MX, Inches(5.33), CW, Inches(0.5),
         "โครงสร้างพื้นฐาน: เซิร์ฟเวอร์ภายในหน่วยงาน │ AI ในเครื่อง │ ฐานข้อมูล │ เครือข่ายปลอดภัย │ ห้ามใช้ AI ภายนอก", 12, WHITE, False, PP_ALIGN.CENTER)

    _txt(s, MX, Inches(6.2), CW, Inches(0.7),
         "ทำไมเริ่มจากระบบการเงิน?\n"
         "1) คดีการเงินมากที่สุด  2) ข้อมูลมีโครงสร้างชัดเจน  3) เห็นผลเร็ว — 3 วัน→30 นาที  4) พื้นฐานสำหรับโมดูลอื่น  5) ไม่มีเครื่องมือฟรีที่ดีในตลาด",
         11, SUBTLE, False, PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════════
    # 3. PROBLEM
    # ════════════════════════════════════════════════════════════
    s = S(LGRAY)
    _title(s, "ปัญหาการสืบสวนทางการเงินในปัจจุบัน")

    ch = Inches(5.6)
    _rect(s, MX, BODY_TOP, hw, ch, WHITE, RED)
    _card_title(s, MX, BODY_TOP, hw, "ปัญหาที่พบ", RED)
    _bullets(s, MX+PAD, BODY_TOP+Inches(0.5), hw-PAD*2, ch-Inches(0.6), [
        "วิเคราะห์รายการเดินบัญชี**ด้วยมือ** ใช้เวลา **3-5 วัน**ต่อคดี",
        "ดูได้**ทีละบัญชี** ไม่เห็น**ภาพรวมเครือข่าย**",
        "ใช้ Excel ธรรมดา ไม่มีเครื่องมือ**วิเคราะห์เฉพาะทาง**",
        "โปรแกรม i2 ราคา **500,000+ บาท/สิทธิ์/ปี** ไม่รองรับภาษาไทย",
        "ไม่มีมาตรฐาน**รายงานส่ง ปปง.**",
        "ไม่มีระบบ**แจ้งเตือนอัตโนมัติ** — พลาดได้ง่าย",
        "ข้อมูลกระจาย**หลาย Excel** ไม่เชื่อมโยงข้ามบัญชี",
        "**ไม่มีเครื่องมือฟรี**ที่ดีสำหรับตำรวจไทย",
    ], bc=RED)

    rx = MX + hw + GAP
    _rect(s, rx, BODY_TOP, hw, ch, WHITE, GREEN)
    _card_title(s, rx, BODY_TOP, hw, "BSIE แก้ได้", GREEN)
    _bullets(s, rx+PAD, BODY_TOP+Inches(0.5), hw-PAD*2, ch-Inches(0.6), [
        "ประมวลผลอัตโนมัติ **ภายใน 30 นาที**",
        "เชื่อมโยงข้ามบัญชี **เห็นเครือข่าย**ทั้งหมด",
        "แผนภาพเครือข่าย แสดงลำดับเวลา **แจ้งเตือน 12 รูปแบบ**",
        "มี**แผนภาพเครือข่ายแบบ i2** ในตัว — **ฟรี**",
        "สร้าง**รายงานพร้อมส่ง ปปง.** ตามมาตรฐาน",
        "**ตรวจจับ 5 รูปแบบฟอกเงิน**ตามมาตรฐานสากลอัตโนมัติ",
        "วิเคราะห์ทางสถิติ: **ค่าเบี่ยงเบน, ค่าผิดปกติ, กฎตัวเลขหลักแรก**",
        "ฟรี — **ไม่มีค่าสิทธิ์รายปี** รองรับ**ทั้งประเทศ**",
    ])

    # ════════════════════════════════════════════════════════════
    # 4. ARCHITECTURE — System
    # ════════════════════════════════════════════════════════════
    s = S(VDARK)
    _title(s, "โครงสร้างระบบ BSIE v4.0", WHITE)

    layer_h = Inches(1.05)
    layer_gap = Inches(0.15)
    y0 = BODY_TOP

    # Frontend
    _rect(s, MX, y0, CW, layer_h, RGBColor(0x20,0x40,0x70))
    _txt(s, MX+PAD, y0+Inches(0.05), Inches(1.5), Inches(0.3), "ส่วนแสดงผล (หน้าจอ)", 14, GOLD, True)
    fe = ["React 19", "TypeScript", "Vite", "Tailwind CSS", "Zustand", "Cytoscape.js", "Recharts", "react-i18next"]
    for i, name in enumerate(fe):
        x = MX + Inches(1.8) + i * Inches(1.3)
        _rect(s, x, y0+Inches(0.55), Inches(1.15), Inches(0.35), DGRAY)
        _txt(s, x, y0+Inches(0.56), Inches(1.15), Inches(0.33), name, 9, SKY, True, PP_ALIGN.CENTER)

    # Backend
    y1 = y0 + layer_h + layer_gap
    _rect(s, MX, y1, CW, layer_h, RGBColor(0x15,0x35,0x25))
    _txt(s, MX+PAD, y1+Inches(0.05), Inches(1.5), Inches(0.3), "ส่วนประมวลผล (เซิร์ฟเวอร์)", 14, MINT, True)
    _txt(s, MX+Inches(1.8), y1+Inches(0.08), Inches(8), Inches(0.3),
         "Python 3.12 │ FastAPI │ 21 ชุดคำสั่ง │ 34 บริการ │ 120+ จุดเชื่อมต่อ", 12, WHITE)
    be = ["ingestion", "search", "graph", "alerts", "fund_flow", "analytics", "reports", "auth", "exports", "workspace"]
    for i, name in enumerate(be):
        x = MX + Inches(1.8) + i * Inches(1.03)
        _rect(s, x, y1+Inches(0.55), Inches(0.9), Inches(0.33), DGRAY)
        _txt(s, x, y1+Inches(0.56), Inches(0.9), Inches(0.31), name, 8, MINT, False, PP_ALIGN.CENTER)

    # Core + Services
    y2 = y1 + layer_h + layer_gap
    _rect(s, MX, y2, hw, layer_h, RGBColor(0x35,0x25,0x15))
    _txt(s, MX+PAD, y2+Inches(0.05), hw-PAD*2, Inches(0.3), "แกนหลักประมวลผล — 30 ชุดโปรแกรม", 13, GOLD, True)
    _txt(s, MX+PAD, y2+Inches(0.38), hw-PAD*2, Inches(0.6),
         "ตรวจจับธนาคาร, จับคู่คอลัมน์, วิเคราะห์ภาษา,\n"
         "จัดรูปแบบข้อมูล, จำแนกประเภท, สร้างความเชื่อมโยง,\n"
         "อ่าน PDF/รูปภาพ, ส่งออกรายงาน, AI ช่วยวิเคราะห์ ...", 10, WHITE)

    _rect(s, MX+hw+GAP, y2, hw, layer_h, RGBColor(0x25,0x15,0x35))
    _txt(s, MX+hw+GAP+PAD, y2+Inches(0.05), hw-PAD*2, Inches(0.3), "34 บริการย่อย", 13, LILAC, True)
    _txt(s, MX+hw+GAP+PAD, y2+Inches(0.38), hw-PAD*2, Inches(0.6),
         "แจ้งเตือน, ตรวจจับผิดปกติ, วิเคราะห์เครือข่าย,\n"
         "ล่าภัยคุกคาม, ติดตามเส้นทางเงิน, ยืนยันตัวตน,\n"
         "รายงาน ปปง., สรุปผล, คิวงาน, ตรวจสอบย้อนหลัง ...", 10, WHITE)

    # Database + AI
    y3 = y2 + layer_h + layer_gap
    _rect(s, MX, y3, hw, Inches(0.8), RGBColor(0x15,0x25,0x15))
    _txt(s, MX+PAD, y3+Inches(0.05), hw-PAD*2, Inches(0.3), "ฐานข้อมูล: 20 ตาราง พร้อมอัปเกรด", 13, MINT, True)
    _txt(s, MX+PAD, y3+Inches(0.38), hw-PAD*2, Inches(0.35),
         "20 ตาราง │ ระบบอัปเกรดอัตโนมัติ │ ปรับแต่งความเร็ว │ อัปเกรดเป็นฐานข้อมูลใหญ่ (เฟส 2)", 10, WHITE)

    _rect(s, MX+hw+GAP, y3, hw, Inches(0.8), RGBColor(0x25,0x15,0x15))
    _txt(s, MX+hw+GAP+PAD, y3+Inches(0.05), hw-PAD*2, Inches(0.3), "ปัญญาประดิษฐ์: AI ภายในเครื่อง", 13, SALMON, True)
    _txt(s, MX+hw+GAP+PAD, y3+Inches(0.38), hw-PAD*2, Inches(0.35),
         "อ่านภาพ/สแกน (ไทย+อังกฤษ) │ ดึงชื่อ/เบอร์/บัตร ปชช./พร้อมเพย์อัตโนมัติ", 10, WHITE)

    # ════════════════════════════════════════════════════════════
    # 5. DATABASE ERD
    # ════════════════════════════════════════════════════════════
    s = S(VDARK)
    _title(s, "ฐานข้อมูล — 20 ตาราง", WHITE)

    tables = [
        ("FileRecord", "ไฟล์ที่อัพโหลด + ตรวจสอบความถูกต้อง", GREEN),
        ("ParserRun", "ประวัติการประมวลผล", GREEN),
        ("Account", "บัญชีทุกแห่ง (จัดรูปแบบแล้ว)", BLUE),
        ("Transaction", "ธุรกรรม 30+ ข้อมูล", BLUE),
        ("StatementBatch", "ชุดรายการเดินบัญชี", GREEN),
        ("RawImportRow", "ข้อมูลดิบสำหรับนิติวิทยาศาสตร์", GREEN),
        ("Entity", "บุคคล/นิติบุคคล", PURPLE),
        ("AccountEntityLink", "เชื่อมบัญชี↔บุคคล", PURPLE),
        ("TransactionMatch", "จับคู่ข้ามบัญชี", ORANGE),
        ("Alert", "แจ้งเตือน + ระดับความรุนแรง", RED),
        ("DuplicateGroup", "กลุ่มข้อมูลซ้ำ", ORANGE),
        ("ReviewDecision", "การตัดสินใจของนักวิเคราะห์", GOLD),
        ("AuditLog", "ร่องรอยตรวจสอบทุกการกระทำ", GOLD),
        ("ExportJob", "งานส่งออกรายงาน", SKY),
        ("User", "ผู้ใช้งาน + ระดับสิทธิ์", SALMON),
        ("AdminSetting", "ตั้งค่า + พื้นที่ทำงาน", TXTL),
        ("GraphAnnotation", "บันทึกบนจุดกราฟ", LILAC),
        ("CaseTag", "แท็กคดี", MINT),
        ("CaseTagLink", "เชื่อมแท็ก↔ข้อมูล", MINT),
        ("MappingProfile", "แบบแผนจับคู่คอลัมน์", GREEN),
    ]
    cw_db = (CW - GAP * 3) / 4
    for i, (name, desc, clr) in enumerate(tables):
        col, row = i % 4, i // 4
        x = MX + col * (cw_db + GAP)
        y = BODY_TOP + row * Inches(1.05)
        _rect(s, x, y, cw_db, Inches(0.85), DGRAY, clr)
        _txt(s, x + Inches(0.12), y + Inches(0.08), cw_db - Inches(0.24), Inches(0.3), name, 12, clr, True)
        _txt(s, x + Inches(0.12), y + Inches(0.42), cw_db - Inches(0.24), Inches(0.35), desc, 10, WHITE)

    _txt(s, MX, Inches(6.5), CW, Inches(0.3),
         "25+ ดัชนีค้นหาเร็ว │ ป้องกันข้อมูลซ้ำ │ เชื่อมโยงตาราง │ ข้อมูลยืดหยุ่น │ อัปเกรดอัตโนมัติ",
         11, SUBTLE, False, PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════════
    # 6. PIPELINE 14 STEPS
    # ════════════════════════════════════════════════════════════
    s = S(VDARK)
    _title(s, "ขั้นตอนการประมวลผล — 14 ขั้นตอน", WHITE)

    steps = [
        ("1", "โหลดไฟล์", "ทุกรูปแบบไฟล์", GREEN),
        ("2", "ตรวจจับธนาคาร", "8 ธนาคารอัตโนมัติ", GREEN),
        ("3", "จับคู่คอลัมน์", "จับคู่อัจฉริยะ 3 ชั้น", BLUE),
        ("4", "โหลดแบบแผน", "จำรูปแบบที่เคยใช้", BLUE),
        ("5", "จัดรูปข้อมูล", "แปลงให้ตรงรูปแบบ", SKY),
        ("6", "ปรับมาตรฐาน", "วันที่/จำนวนเงิน/เข้า-ออก", SKY),
        ("7", "แยกเลขบัญชี", "แยกและจัดรูปแบบ", ORANGE),
        ("8", "แยกรายละเอียด", "แยกรายละเอียด", ORANGE),
        ("9", "ดึงชื่อ/เบอร์/ID", "ชื่อไทย/เบอร์/บัตร ปชช.", PURPLE),
        ("10", "จำแนกประเภท", "เงินเข้า/เงินออก/โอน", PURPLE),
        ("11", "สร้างเชื่อมโยง", "จากบัญชี→ไปบัญชี", GOLD),
        ("12", "แก้ไขด้วยมือ", "แก้ไขจากผู้วิเคราะห์", GOLD),
        ("13", "สร้างทะเบียนบุคคล", "สร้าง Entity list", SALMON),
        ("14", "ส่งออกรายงาน", "ชุดข้อมูลบัญชี", SALMON),
    ]
    sw_step = (CW - GAP * 6) / 7
    for i, (num, name, desc, clr) in enumerate(steps):
        col, row = i % 7, i // 7
        x = MX + col * (sw_step + GAP)
        y = BODY_TOP + row * Inches(2.85)
        sh = Inches(2.5)
        _rect(s, x, y, sw_step, sh, DGRAY, clr)
        _txt(s, x, y+Inches(0.15), sw_step, Inches(0.5), num, 26, clr, True, PP_ALIGN.CENTER)
        _txt(s, x, y+Inches(0.7), sw_step, Inches(0.35), name, 12, WHITE, True, PP_ALIGN.CENTER)
        _txt(s, x+Inches(0.08), y+Inches(1.15), sw_step-Inches(0.16), Inches(1.0), desc, 10, RGBColor(0xBE,0xE3,0xF8), False, PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════════
    # 7. DATA INTAKE
    # ════════════════════════════════════════════════════════════
    s = S(LGRAY)
    _title(s, "การนำเข้าข้อมูล — รองรับทุกรูปแบบ")

    intake = [
        ("8 ธนาคารไทย Auto-detect", [
            "**SCB** — ไทยพาณิชย์", "**KBANK** — กสิกรไทย",
            "**BBL** — กรุงเทพ", "**KTB** — กรุงไทย",
            "**BAY** — กรุงศรี", "**TTB** — ทหารไทยธนชาต",
            "**GSB** — ออมสิน", "**BAAC** — ธ.ก.ส.",
        ]),
        ("รูปแบบไฟล์ + OCR", [
            "**Excel** (.xlsx / .xls) — ทุกธนาคาร",
            "**PDF** ข้อความ — pdfplumber extraction",
            "**PDF สแกน** — EasyOCR + table reconstruction",
            "**รูปภาพ** (JPG/PNG/BMP) — EasyOCR Thai+EN",
            "**OFX** — ธนาคารออนไลน์",
            "**SHA-256** ตรวจจับไฟล์ซ้ำอัตโนมัติ",
        ]),
        ("จับคู่คอลัมน์อัจฉริยะ", [
            "**จับคู่อัจฉริยะ 3 ชั้น** — ตรง → ชื่อคล้าย → คล้ายคลึง",
            "**เรียนรู้เอง** — จำรูปแบบที่เคยใช้",
            "**ตรวจเนื้อหา** — อ่านข้อความในไฟล์ช่วยตัดสิน",
            "**ทำใหม่ได้** — จับคู่ผิด? สั่งประมวลผลซ้ำ",
            "**ดึงข้อมูลอัตโนมัติ** — ชื่อไทย, เบอร์โทร, บัตร ปชช., พร้อมเพย์",
        ]),
    ]
    for i, (t_, items) in enumerate(intake):
        x = MX + i * (tw + GAP)
        _rect(s, x, BODY_TOP, tw, Inches(5.8), WHITE)
        _card_title(s, x, BODY_TOP, tw, t_)
        _bullets(s, x+PAD, BODY_TOP+Inches(0.55), tw-PAD*2, Inches(5.0), items, sz=12)

    # ════════════════════════════════════════════════════════════
    # 8. LINK CHART
    # ════════════════════════════════════════════════════════════
    s = S(NAVY)
    _title(s, "แผนภาพเครือข่ายการเงิน — แบบ i2", WHITE)
    _sub(s, "แผนภาพเชื่อมโยงแบบโต้ตอบ — กดขยายได้ไม่จำกัด", RGBColor(0xBE,0xE3,0xF8))

    ch = Inches(5.1)
    _rect(s, MX, Inches(1.3), hw, ch, CDARK)
    _card_title(s, MX, Inches(1.3), hw, "แผนภาพเชื่อมโยงแบบโต้ตอบ", GOLD)
    _bullets(s, MX+PAD, Inches(1.85), hw-PAD*2, ch-Inches(0.7), [
        "กดที่บัญชีใดก็ได้ **ขยายเห็นเครือข่ายต่อได้เรื่อยๆ**",
        "**5 รูปแบบการจัด**: วงกลม, กระจาย, กระชับ, ลำดับชั้น, พัด",
        "**ขนาดจุดตามยอดเงิน** — ยอดมาก จุดใหญ่",
        "**สีขอบตามระดับเสี่ยง** — แดง = น่าสงสัยมาก",
        "**แสดงยอดเงิน**บนเส้นเชื่อมได้",
        "**ปักหมุด / ซ่อน / เลือกหลายจุด** จัดระเบียบได้",
        "**ประวัติการค้นหา** — ย้อนดูจุดที่เคยดูได้",
        "**บันทึกเป็นรูปภาพ** ใส่รายงานได้ทันที",
    ], clr=WHITE, bc=GOLD, sz=12)

    rx = MX + hw + GAP
    _rect(s, rx, Inches(1.3), hw, ch, CDARK)
    _card_title(s, rx, Inches(1.3), hw, "วิเคราะห์เครือข่าย + ติดตามเส้นทางเงิน", GOLD)
    _bullets(s, rx+PAD, Inches(1.85), hw-PAD*2, ch-Inches(0.7), [
        "**บัญชีที่เชื่อมต่อมากที่สุด** — ศูนย์กลางเครือข่าย",
        "**บัญชีตัวกลาง** — คนกลางที่เงินผ่านบ่อย",
        "**บัญชีที่เข้าถึงทุกคนเร็วที่สุด**",
        "**ติดตามเส้นทางเงิน** — ตามได้สูงสุด **4 ทอด**",
        "**จดบันทึก** — เขียนโน้ตบนจุดบัญชีได้",
        "**ติดป้าย**: ผู้ต้องสงสัย, พยาน, เหยื่อ, ตรวจสอบเพิ่ม",
        "**พื้นที่ทำงาน** — บันทึกแผนภาพไว้ดูทีหลังได้",
        "**ประวัติบุคคล** — หน้ารวมข้อมูลทุกบัญชีของคนเดียวกัน",
    ], clr=WHITE, bc=GOLD, sz=12)

    # ════════════════════════════════════════════════════════════
    # 9. VISUALIZATION
    # ════════════════════════════════════════════════════════════
    s = S(LGRAY)
    _title(s, "การแสดงผลข้อมูล — กราฟเวลา, แผนที่ความร้อน, หน้าสรุป")

    viz = [
        ("กราฟเวลา", BLUE, [
            "**แท่งกราฟ + จุด** แสดงธุรกรรมตามวัน",
            "เลือกดูได้: **รายวัน / รายสัปดาห์ / รายเดือน**",
            "ดูรูปแบบธุรกรรมตามเวลา",
            "**พบช่วงเวลาที่มีธุรกรรมผิดปกติ**",
        ]),
        ("แผนที่ความร้อน", PURPLE, [
            "**แสดงผลตาม ชั่วโมง × วัน**",
            "ครอบคลุม **24 ชม. × 7 วัน = 168 ช่อง**",
            "**สีเข้ม = ธุรกรรมถี่** สีอ่อน = น้อย",
            "**ตรวจจับเวลาผิดปกติ** เช่น ตี 2-5",
        ]),
        ("หน้าสรุปภาพรวม", GREEN, [
            "**สรุปภาพรวม** — จำนวนเงินเข้า/ออก, สถิติ",
            "**ธุรกรรมล่าสุด** — ดูรายการใหม่ได้ทันที",
            "**บัญชียอดสูงสุด** — เรียงตามจำนวนเงิน",
            "**โหลดเร็ว** — แสดงผลเฉพาะส่วนที่ต้องการ",
        ]),
    ]
    for i, (t_, clr, items) in enumerate(viz):
        x = MX + i * (tw + GAP)
        _rect(s, x, BODY_TOP, tw, Inches(4.5), WHITE, clr)
        _card_title(s, x, BODY_TOP, tw, t_, clr)
        _bullets(s, x+PAD, BODY_TOP+Inches(0.55), tw-PAD*2, Inches(3.5), items)

    # ════════════════════════════════════════════════════════════
    # 10. ALERTS
    # ════════════════════════════════════════════════════════════
    s = S(ORANGE)
    _title(s, "ระบบแจ้งเตือนอัตโนมัติ — 12 รูปแบบ", WHITE)

    alert_data = [
        ("7 กฎตรวจจับจากเครือข่าย", [
            "**โอนซ้ำ** — บัญชีคู่เดิมโอนกัน 3+ ครั้ง",
            "**เงินเข้าหลายทาง** — รับจาก 3+ บัญชี",
            "**เงินออกหลายทาง** — โอนไป 3+ บัญชี",
            "**เงินวนกลับ** — A โอนให้ B แล้ว B โอนกลับ",
            "**บัญชีผ่านมือ** — รับแล้วโอนต่อทันที",
            "**บัญชีศูนย์กลาง** — เชื่อมกับ 6+ บัญชี",
            "**คู่ค้าซ้ำ** — เจอคนเดิมหลายวัน",
        ]),
        ("5 รูปแบบฟอกเงินตามมาตรฐานสากล", [
            "**แบ่งรายการย่อย** — โอนทีละน้อยหลายครั้ง หลีกเลี่ยงเกณฑ์",
            "**สร้างชั้นซ้อน** — โอนผ่านหลายบัญชีเพื่อปิดบัง",
            "**เงินเข้า-ออกเร็ว** — เข้ามาแล้วออกภายใน 24 ชม.",
            "**บัญชีนิ่งแล้วตื่น** — ไม่ใช้นานแล้วมีเงินเข้าทันที",
            "**เงินวนกลับ** — ส่งออกไปแล้ววนกลับมาที่เดิม",
        ]),
        ("4 วิธีตรวจจับทางคณิตศาสตร์", [
            "**ค่าเบี่ยงเบน** — ยอดเงินผิดปกติจากค่าเฉลี่ย",
            "**ค่าผิดปกติ** — ยอดเงินที่สูง/ต่ำกว่าพิสัยปกติ",
            "**กฎตัวเลขหลักแรก** — ตัวเลขตั้งต้นไม่เป็นธรรมชาติ",
            "**ค่าเฉลี่ยเคลื่อนที่** — เปรียบเทียบกับ 30 รายการก่อน",
        ]),
    ]
    for i, (t_, items) in enumerate(alert_data):
        x = MX + i * (tw + GAP)
        _rect(s, x, BODY_TOP, tw, Inches(5.6), RGBColor(0x9C,0x4E,0x18))
        _card_title(s, x, BODY_TOP, tw, t_, LGOLD)
        _bullets(s, x+PAD, BODY_TOP+Inches(0.55), tw-PAD*2, Inches(4.8), items, clr=WHITE, bc=LGOLD, sz=12, sp=Pt(3))

    # ════════════════════════════════════════════════════════════
    # 11. CROSS-ACCOUNT + REPORTS
    # ════════════════════════════════════════════════════════════
    s = S(LGRAY)
    _title(s, "การวิเคราะห์ข้ามบัญชี + โต๊ะทำงานสืบสวน + รายงาน")

    _rect(s, MX, BODY_TOP, hw, Inches(2.7), WHITE, BLUE)
    _card_title(s, MX, BODY_TOP, hw, "วิเคราะห์ข้ามบัญชี", BLUE)
    _bullets(s, MX+PAD, BODY_TOP+Inches(0.5), hw-PAD*2, Inches(2.0), [
        "**ดูเส้นทางเงิน** — ดูเงินเข้า/ออกของบัญชีใดก็ได้",
        "**เปรียบเทียบคู่** — ดูธุรกรรมทั้งหมดระหว่าง 2 บัญชี",
        "**ติดตามเส้นทาง** — ตามเงินจาก A ไป B ไป C ได้ถึง **4 ทอด**",
        "**จับคู่ข้ามบัญชี** — ตรวจทุกบัญชีพร้อมกันในครั้งเดียว",
        "**เปรียบเทียบช่วงเวลา** — เทียบ 2 ช่วงเวลาหาความต่าง",
    ], sz=12)

    rx = MX + hw + GAP
    _rect(s, rx, BODY_TOP, hw, Inches(2.7), WHITE, PURPLE)
    _card_title(s, rx, BODY_TOP, hw, "โต๊ะทำงานสืบสวน — 13 เครื่องมือ", PURPLE)
    _bullets(s, rx+PAD, BODY_TOP+Inches(0.5), hw-PAD*2, Inches(2.0), [
        "ฐานข้อมูล │ ไฟล์ │ ประมวลผล │ **บัญชี** │ **ค้นหา**",
        "**แจ้งเตือน** │ **ข้ามบัญชี** │ **แผนภาพเชื่อมโยง** │ **ลำดับเวลา**",
        "ข้อมูลซ้ำ │ จับคู่ │ **ตรวจสอบย้อนหลัง** │ **ส่งออก**",
    ], sz=12, bc=PURPLE)

    # Reports
    _rect(s, MX, BODY_TOP+Inches(2.95), CW, Inches(3.1), WHITE, GOLD)
    _card_title(s, MX, BODY_TOP+Inches(2.95), CW, "รายงานพร้อมใช้ — 6 รูปแบบ", GOLD)
    reports = [
        ("Excel (.xlsx)", "14+ ชีท, TH Sarabun New, สีตามประเภท, สูตร =SUM/=COUNTA, auto-filter"),
        ("PDF", "ปก + สถิติ + คู่สัญญา + alerts + ธุรกรรม + ช่องลงนาม"),
        ("แผนภาพ i2 (.anx)", "เปิดใน Analyst's Notebook ได้ทันที — XML format"),
        ("นำเข้า i2 (.csv+.ximp)", "ข้อมูล 29 คอลัมน์ + ไฟล์นำเข้า i2"),
        ("STR / CTR", "รูปแบบรายงาน ปปง. พร้อมส่ง — regulatory compliance"),
        ("CSV", "transactions, entities, links, reconciliation, graph data"),
    ]
    for i, (fmt, desc) in enumerate(reports):
        col, row = i % 2, i // 2
        x = MX + PAD + col * (CW / 2)
        y = BODY_TOP + Inches(3.55) + row * Inches(0.55)
        _txt(s, x, y, Inches(1.8), Inches(0.5), fmt, 12, NAVY, True)
        _txt(s, x + Inches(1.8), y, Inches(4.0), Inches(0.5), desc, 10, TXTL)

    # ════════════════════════════════════════════════════════════
    # 12. SECURITY
    # ════════════════════════════════════════════════════════════
    s = S(VDARK)
    _title(s, "ความปลอดภัย, สายโซ่หลักฐาน, ภาษาไทย", WHITE)

    sec = [
        ("ยืนยันตัวตน", SKY, [
            "**ล็อกอินด้วยรหัส** + 3 ระดับ: ผู้ดูแล/นักวิเคราะห์/ผู้ชม",
            "**เข้ารหัสรหัสผ่าน** 100,000 รอบ (ปลอดภัยสูง)",
            "**จำกัดการล็อกอินผิด** — ป้องกันเดารหัส",
            "**ตั้งค่าได้** — เปิด/ปิดระบบยืนยันตัวตน",
        ]),
        ("ป้องกันการโจมตี", SALMON, [
            "**ป้องกันฝังในเว็บอื่น**",
            "**ป้องกันปลอมแปลงไฟล์**",
            "**ไม่เปิดเผยที่มา**",
            "**จำกัดประเภทไฟล์** — เฉพาะ Excel/PDF/CSV/รูปภาพ",
            "**จำกัดขนาดไฟล์** — ไม่เกิน 50 MB",
        ]),
        ("สายโซ่หลักฐาน", GOLD, [
            "**ร่องรอยตรวจสอบ** — บันทึกว่าใคร/ทำอะไร/เมื่อไหร่",
            "**สายโซ่หลักฐาน** — ประวัติครบถ้วนใช้เป็นหลักฐานได้",
            "**เรียกดูประวัติ** ผ่านระบบได้ทุกเมื่อ",
            "**ตรวจสอบความถูกต้อง** ของไฟล์หลักฐาน",
            "**ตรวจข้อมูลไฟล์** — ใช้ทางนิติวิทยาศาสตร์",
        ]),
    ]
    for i, (t_, clr, items) in enumerate(sec):
        x = MX + i * (tw + GAP)
        _rect(s, x, BODY_TOP, tw, Inches(4.3), DGRAY)
        _card_title(s, x, BODY_TOP, tw, t_, clr)
        _bullets(s, x+PAD, BODY_TOP+Inches(0.5), tw-PAD*2, Inches(3.5), items, sz=11, clr=WHITE, bc=clr, sp=Pt(2))

    # Thai bar
    _rect(s, MX, Inches(5.7), CW, Inches(0.8), DGRAY, GOLD)
    _txt(s, MX+PAD, Inches(5.75), CW-PAD*2, Inches(0.7),
         "ภาษาไทยเต็มระบบ: สลับ ไทย/อังกฤษ ได้ทั้งระบบ │ แปลแล้วกว่า 500 จุด │ ฟอนต์ TH Sarabun New ทุกรายงาน\n"
         "ดึงชื่อไทย/เบอร์โทร/บัตร ปชช./พร้อมเพย์อัตโนมัติ │ ทดสอบแล้ว 244 ชุด",
         11, WHITE, False, PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════════
    # 13. LOCAL LLM
    # ════════════════════════════════════════════════════════════
    s = S(DRED)
    _title(s, "ทำไมต้องใช้ AI ภายในเครื่อง — ห้ามใช้ AI ภายนอก", WHITE)

    ch = Inches(5.3)
    _rect(s, MX, BODY_TOP, hw, ch, RGBColor(0x74,0x2A,0x2A))
    _card_title(s, MX, BODY_TOP, hw, "เหตุผลความจำเป็น", SALMON)
    _bullets(s, MX+PAD, BODY_TOP+Inches(0.5), hw-PAD*2, ch-Inches(0.7), [
        "**ความลับทางราชการ** — ระดับ 'ลับ' / 'ลับมาก'",
        "**พ.ร.บ. คุ้มครองข้อมูลส่วนบุคคล**",
        "**เสี่ยงถูกนำข้อมูลไปเรียนรู้** — AI ภายนอกอาจเอาข้อมูลคดีไปใช้",
        "**สายโซ่หลักฐาน** — หลักฐานต้องไม่ผ่านระบบภายนอก",
        "**ระเบียบ สตช.** — ข้อมูลสืบสวนต้องอยู่ในเครือข่ายภายใน",
        "**ป้องกันข้อมูลรั่วไหล** — เลขบัญชี, ชื่อผู้ต้องสงสัย, เส้นทางเงิน",
    ], clr=WHITE, bc=SALMON, sz=13)

    rx = MX + hw + GAP
    _rect(s, rx, BODY_TOP, hw, ch, RGBColor(0x74,0x2A,0x2A))
    _card_title(s, rx, BODY_TOP, hw, "AI ภายในเครื่องทำอะไรได้", SALMON)
    _bullets(s, rx+PAD, BODY_TOP+Inches(0.5), hw-PAD*2, Inches(3.3), [
        "**สรุปคดีให้อัตโนมัติ** — เช่น 'พบเงินเข้าจาก 15 บัญชี รวม 3.2 ล้านบาท'",
        "**ตอบคำถาม** — เช่น 'บัญชีนี้มีธุรกรรมเกิน 1 แสนกี่รายการ?'",
        "**จำแนกธุรกรรม** — ช่วยตัดสินรายการที่ระบบไม่มั่นใจ",
        "**ตรวจจับรูปแบบ** — มีการแบ่งรายการย่อยหรือไม่?",
        "**ร่างรายงาน ปปง.** ให้อัตโนมัติ",
    ], clr=WHITE, bc=SALMON, sz=13)
    _txt(s, rx+PAD, BODY_TOP+Inches(4.1), hw-PAD*2, Inches(1.0),
         "แนะนำ: ใช้โปรแกรม Ollama + โมเดล Llama 3.1 หรือ Typhoon (เข้าใจภาษาไทย)\n"
         "ต้องมีการ์ดจอ: NVIDIA RTX 4090 หรือ A6000\n"
         "ติดตั้งง่าย — ใช้งานได้ทันที", 11, RGBColor(0xFE,0xBD,0xBD))

    # ════════════════════════════════════════════════════════════
    # 14. INFRASTRUCTURE
    # ════════════════════════════════════════════════════════════
    s = S(VDARK)
    _title(s, "โครงสร้างพื้นฐาน — เฟส 1 ถึง 4", WHITE)

    inf = [
        ("เฟส 1: เครื่องเดี่ยว (ฟรี)", MINT,
         "ใช้คอมพิวเตอร์ที่มี → เปิดระบบ BSIE\n"
         "ทำงานบนเครื่องเดียว ไม่ต้องต่อเน็ต\n"
         "ไม่ต้องซื้อเซิร์ฟเวอร์"),
        ("เฟส 2: ใช้ร่วมกันทั้งจังหวัด", SKY,
         "ตั้งเซิร์ฟเวอร์กลางจังหวัด ปลอดภัยด้วยรหัสลับ\n"
         "ฐานข้อมูลใหญ่ + AI ภายในเครื่อง + การ์ดจอ\n"
         "พนักงานสอบสวน 1-30 คน → เปิดผ่านเบราว์เซอร์"),
    ]
    for i, (t_, clr, desc) in enumerate(inf):
        x = MX + i * (hw + GAP)
        _rect(s, x, BODY_TOP, hw, Inches(2.6), DGRAY, clr)
        _card_title(s, x, BODY_TOP, hw, t_, clr)
        _txt(s, x+PAD, BODY_TOP+Inches(0.5), hw-PAD*2, Inches(1.9), desc, 11, WHITE)

    _rect(s, MX, BODY_TOP+Inches(2.85), CW, Inches(3.0), DGRAY, GOLD)
    _card_title(s, MX, BODY_TOP+Inches(2.85), CW, "เฟส 3-4: ระดับจังหวัด / ภาค / ประเทศ (เซิร์ฟเวอร์ภายในหน่วยงาน)", GOLD)
    _txt(s, MX+PAD, BODY_TOP+Inches(3.4), hw-PAD, Inches(2.2),
         "Data Center จังหวัด / ภาค\n"
         "  ├── Load Balancer (Nginx)\n"
         "  ├── BSIE Backend ×2-3\n"
         "  ├── PostgreSQL (primary+replica)\n"
         "  ├── Redis (cache + session)\n"
         "  ├── NAS/SAN (evidence)\n"
         "  ├── Monitoring (Prometheus+Grafana)\n"
         "  └── Backup server", 11, WHITE)
    _txt(s, MX+hw+GAP, BODY_TOP+Inches(3.4), hw-PAD, Inches(2.2),
         "Local LLM Cluster\n"
         "  ├── GPU Server ×1-2 per ภาค\n"
         "  ├── NVIDIA A6000 ×2 per server\n"
         "  └── Llama 3.1 70B + Typhoon\n\n"
         "เชื่อม VPN/HTTPS:\n"
         "  สถานี อ.เมือง → VPN\n"
         "  สถานี อ.เกาะสมุย → VPN\n"
         "  บก.จว. → VPN", 11, WHITE)

    # ════════════════════════════════════════════════════════════
    # 15. ROADMAP — Horizontal Timeline
    # ════════════════════════════════════════════════════════════
    s = S(VDARK)
    _title(s, "แผนพัฒนาระบบ SPNI", WHITE)

    TL_Y = Inches(2.55)  # timeline center

    # Main horizontal axis
    _line(s, MX, TL_Y, CW, Pt(3), RGBColor(0x4A,0x55,0x68))

    # Year marks
    for label, x in [("พ.ศ. 2569", MX), ("2570", MX+Inches(3.9)),
                      ("2571", MX+Inches(7.8)), ("2572+", MX+Inches(11.2))]:
        _line(s, x, TL_Y-Inches(0.15), Pt(2), Inches(0.3), RGBColor(0x71,0x81,0x96))
        _txt(s, x-Inches(0.4), TL_Y-Inches(0.45), Inches(0.9), Inches(0.25),
             label, 10, SUBTLE, True, PP_ALIGN.CENTER)

    # Phase bars
    phases_tl = [
        ("เฟส 1\nทดลอง", "เดือน 1-3", MX, Inches(2.3), GREEN, "ฟรี",
         ["BSIE v4.0 Pilot", "ทดลอง 5-10 คดีจริง", "เก็บ feedback"]),
        ("เฟส 2\nจังหวัดนำร่อง", "เดือน 4-9", MX+Inches(2.5), Inches(2.8), BLUE, "~8 ล้าน",
         ["Server 5 จังหวัด", "PostgreSQL + Multi-user", "Local LLM + GPU", "อบรม 5 จังหวัด"]),
        ("เฟส 3\nขยายระบบ", "เดือน 10-21", MX+Inches(5.5), Inches(3.8), ORANGE, "~40 ล้าน",
         ["CDR Module (โมดูล 2)", "Shared Intelligence Layer", "Mobile App", "ขยาย 20 จังหวัด"]),
        ("Phase 4\nระดับประเทศ", "ปีที่ 3+", MX+Inches(9.5), Inches(2.63), RED, "~120 ล้าน",
         ["Social + CCTV Module", "77 จว. ทั้งประเทศ", "Data Center ×10 ภาค", "เชื่อม ปปง./DSI"]),
    ]
    for label, time, x, w, clr, budget, items in phases_tl:
        # Phase bar
        _rect(s, x, TL_Y - Inches(0.28), w, Inches(0.55), clr)
        _txt(s, x, TL_Y - Inches(0.28), w, Inches(0.32), label.split('\n')[0], 11, WHITE, True, PP_ALIGN.CENTER)
        if '\n' in label:
            _txt(s, x, TL_Y - Inches(0.02), w, Inches(0.25), label.split('\n')[1], 9, WHITE, False, PP_ALIGN.CENTER)

        # Diamond milestone
        dm = Inches(0.15)
        d = s.shapes.add_shape(MSO_SHAPE.DIAMOND, x - dm/2, TL_Y - dm/2, dm, dm)
        d.fill.solid(); d.fill.fore_color.rgb = clr; d.line.fill.background()

        # Budget badge below timeline
        bw = min(w, Inches(1.6))
        _rect(s, x + (w-bw)/2, TL_Y + Inches(0.2), bw, Inches(0.35), DGRAY, clr)
        _txt(s, x + (w-bw)/2, TL_Y + Inches(0.22), bw, Inches(0.3),
             budget, 11, clr, True, PP_ALIGN.CENTER)

        # Time label
        _txt(s, x, TL_Y + Inches(0.6), w, Inches(0.2), time, 9, SUBTLE, False, PP_ALIGN.CENTER)

        # Deliverables above
        card_h = Inches(0.22 * len(items) + 0.15)
        card_y = TL_Y - Inches(0.4) - card_h
        _rect(s, x, card_y, w, card_h, DGRAY)
        for j, item in enumerate(items):
            _txt(s, x + Inches(0.1), card_y + Inches(0.08 + j * 0.22),
                 w - Inches(0.2), Inches(0.2), "● " + item, 9, WHITE)

    # Grand total bar
    _rect(s, MX, Inches(6.2), CW, Inches(0.55), DGRAY, GOLD)
    _txt(s, MX+PAD, Inches(6.23), Inches(4), Inches(0.5), "งบประมาณรวม 3 ปี", 15, WHITE, True)
    _txt(s, Inches(5), Inches(6.23), CW - Inches(4.5), Inches(0.5),
         "~170 ล้านบาท │ ฟรี → 8 ล้าน → 40 ล้าน → 120 ล้าน", 14, GOLD, True, PP_ALIGN.RIGHT)

    _txt(s, MX, Inches(6.9), CW, Inches(0.3),
         "เฟส 1 เริ่มได้ทันที (ฟรี) → เฟส 2 ของบหลังทดลองสำเร็จ → เฟส 3-4 ขยายตามผลลัพธ์จริง",
         11, SUBTLE, False, PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════════
    # 16. BUDGET NATIONAL
    # ════════════════════════════════════════════════════════════
    s = S(GOLD)
    _title(s, "งบประมาณระดับประเทศ — ระบบ SPNI", VDARK)

    budget_phases = [
        ("เฟส 1: ทดลอง (3 เดือน)", [
            ("ใช้เครื่อง PC ที่มีอยู่ + ซอฟต์แวร์ BSIE v4.0", "ฟรี"),
        ], "ฟรี", GREEN),
        ("เฟส 2: นำร่อง 5 จังหวัด (6 เดือน)", [
            ("เซิร์ฟเวอร์หลัก ×5 (CPU 8-core, 32GB, 1TB SSD)", "1,500,000 ฿"),
            ("เซิร์ฟเวอร์ AI + การ์ดจอ ×5 (RTX 4090/A6000)", "5,000,000 ฿"),
            ("Network / UPS / อุปกรณ์ ×5 จังหวัด", "500,000 ฿"),
            ("อัปเกรดฐานข้อมูล + พัฒนาเพิ่ม", "500,000 ฿"),
            ("อบรมพนักงานสอบสวน 5 จังหวัด", "300,000 ฿"),
        ], "~8 ล้าน", BLUE),
        ("เฟส 3: ขยาย 20 จังหวัด + โมดูล 2 (12 เดือน)", [
            ("ระบบวิเคราะห์โทรศัพท์ (ส่วนงาน 2)", "8,000,000 ฿"),
            ("Server + GPU ×20 จังหวัด", "20,000,000 ฿"),
            ("ระบบข่าวกรองกลาง + แอปมือถือ", "7,000,000 ฿"),
            ("ทีมพัฒนา 5-8 คน × 12 เดือน", "5,000,000 ฿"),
        ], "~40 ล้าน", ORANGE),
        ("เฟส 4: ทั้งประเทศ 77 จังหวัด + 10 ภาค (ปีที่ 2-3)", [
            ("ระบบสื่อสังคม + กล้องวงจรปิด (ส่วนงาน 3-4)", "30,000,000 ฿"),
            ("ศูนย์ข้อมูลระดับภาค ×10 + GPU Cluster", "40,000,000 ฿"),
            ("ติดตั้ง 77 จังหวัด — server + network", "30,000,000 ฿"),
            ("ทีมพัฒนา 15 คน + อบรมทั่วประเทศ", "20,000,000 ฿"),
        ], "~120 ล้าน", RED),
    ]
    y = BODY_TOP - Inches(0.1)
    for phase_name, items, total, clr in budget_phases:
        h = Inches(0.28 * max(len(items), 1) + 0.35)
        _rect(s, MX, y, CW - Inches(2.7), h, RGBColor(0xFF,0xF8,0xE8))
        _txt(s, MX+PAD, y+Inches(0.04), Inches(7), Inches(0.28), phase_name, 12, clr, True)
        for j, (item, cost) in enumerate(items):
            _txt(s, MX+Inches(0.4), y+Inches(0.3+j*0.28), Inches(5.5), Inches(0.25), item, 10, TXT)
            _txt(s, MX+Inches(6.0), y+Inches(0.3+j*0.28), Inches(3.0), Inches(0.25), cost, 10, TXTL, False, PP_ALIGN.RIGHT)
        # Badge
        _rect(s, CW - Inches(2.0), y+Inches(0.08), Inches(2.3), Inches(0.4), clr)
        _txt(s, CW - Inches(2.0), y+Inches(0.08), Inches(2.3), Inches(0.4), total, 14, WHITE, True, PP_ALIGN.CENTER)
        y += h + Inches(0.08)

    _rect(s, MX, y+Inches(0.05), CW, Inches(0.55), VDARK)
    _txt(s, MX+PAD, y+Inches(0.08), Inches(4), Inches(0.45), "รวมทั้งโครงการ 3 ปี", 16, WHITE, True)
    _txt(s, Inches(5.5), y+Inches(0.08), CW-Inches(5.5), Inches(0.45), "~170 ล้านบาท", 22, GOLD, True, PP_ALIGN.RIGHT)

    # ════════════════════════════════════════════════════════════
    # 17. BUDGET HIGHLIGHTS
    # ════════════════════════════════════════════════════════════
    s = S(LGRAY)
    _title(s, "จุดเด่นของงบประมาณ SPNI")

    highlights = [
        ("ไม่มีค่าสิทธิ์ใช้งานรายปี", "ซอฟต์แวร์ทั้งหมดพัฒนาเอง + Open Source — จ่ายครั้งเดียว ใช้ตลอด\nเปรียบเทียบ: i2 = 500,000 ฿/คน/ปี × 1,484 สถานี = 742 ล้าน/ปี", GREEN),
        ("ขยายได้ไม่จำกัด", "เพิ่มสถานีใหม่ = แค่เพิ่ม server ไม่ต้องจ่าย license เพิ่ม\nรองรับ 1,484 สถานี ~15,000-20,000 พนักงานสอบสวนทั่วประเทศ", BLUE),
        ("ปรับแต่งได้ตามต้องการ", "มีโค้ดทั้งหมด — แก้ไข ปรับแต่ง เพิ่มฟีเจอร์ ได้เอง\nไม่ต้องพึ่งบริษัทต่างชาติ — เป็นสมบัติของ สตช.", ORANGE),
        ("อธิปไตยข้อมูล", "ข้อมูลทั้งหมดอยู่ในเซิร์ฟเวอร์ภายในหน่วยงานของ สตช.\nไม่มีข้อมูลคดีออกนอกประเทศ — Local LLM เท่านั้น", RED),
    ]
    for i, (t_, desc, clr) in enumerate(highlights):
        y = BODY_TOP + i * Inches(1.4)
        _line(s, MX, y, Inches(0.12), Inches(1.2), clr)
        _rect(s, MX + Inches(0.2), y, CW - Inches(0.2), Inches(1.2), WHITE)
        _txt(s, MX + Inches(0.4), y + Inches(0.08), CW - Inches(0.8), Inches(0.35), t_, 17, clr, True)
        _txt(s, MX + Inches(0.4), y + Inches(0.48), CW - Inches(0.8), Inches(0.65), desc, 12, TXTL)

    # ════════════════════════════════════════════════════════════
    # 18. COMPARISON
    # ════════════════════════════════════════════════════════════
    s = S(LGRAY)
    _title(s, "เปรียบเทียบกับโปรแกรมที่ขายในตลาด")

    rows = [
        ("คุณสมบัติ", "BSIE (SPNI)", "i2 Analyst's NB", "Cellebrite Fin."),
        ("ราคา", "ฟรี", "500,000 ฿/คน/ปี", "300,000 ฿/คน/ปี"),
        ("ค่าใช้จ่าย 1,484 สถานี/ปี", "ค่า server เท่านั้น", "742+ ล้าน/ปี", "445+ ล้าน/ปี"),
        ("ภาษาไทย", "✓ เต็มระบบ", "✗", "✗"),
        ("ธนาคารไทย 8 แห่ง", "✓ ตรวจจับอัตโนมัติ", "✗ setup เอง", "~ บางส่วน"),
        ("กราฟเครือข่าย", "✓ mini i2 (5 layout)", "✓ เต็มรูปแบบ", "✓"),
        ("PDF/OCR ภาษาไทย", "✓ EasyOCR", "✗", "~"),
        ("รายงาน ปปง.", "✓ STR/CTR", "✗", "~"),
        ("PromptPay", "✓", "✗", "✗"),
        ("Threat Hunting FATF", "✓ 5 รูปแบบ", "✗", "~"),
        ("Anomaly Detection", "✓ 4 วิธี", "~ บางส่วน", "~"),
        ("Local LLM AI", "✓ Ollama", "✗", "✗"),
        ("ปรับแต่ง", "✓ มีโค้ดทั้งหมด", "✗ แก้ไขไม่ได้", "✗ แก้ไขไม่ได้"),
        ("อธิปไตยข้อมูล", "✓ ภายในหน่วย 100%", "~ ขึ้นกับ deploy", "~"),
    ]
    col_ws = [Inches(3.2), Inches(3.0), Inches(3.0), Inches(3.0)]
    col_xs = [MX]
    for w in col_ws[:-1]:
        col_xs.append(col_xs[-1] + w + Inches(0.02))

    rh = Inches(0.38)
    for i, row in enumerate(rows):
        y = BODY_TOP + i * rh
        is_hdr = (i == 0)
        for j, cell in enumerate(row):
            bg_c = NAVY if is_hdr else (WHITE if i % 2 == 1 else MGRAY)
            tc = WHITE if is_hdr else TXT
            if j == 1 and not is_hdr and ("✓" in cell or cell == "ฟรี"):
                tc = GREEN
            elif not is_hdr and "✗" in cell:
                tc = RED
            elif not is_hdr and cell.startswith("~"):
                tc = GOLD
            _rect(s, col_xs[j], y, col_ws[j], rh, bg_c)
            al = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            bld = is_hdr or (j == 0)
            _txt(s, col_xs[j]+Inches(0.08), y+Inches(0.02), col_ws[j]-Inches(0.16), rh-Inches(0.04),
                 cell, 11 if not is_hdr else 12, tc, bld, al)

    # ════════════════════════════════════════════════════════════
    # 19. PROJECT STATS
    # ════════════════════════════════════════════════════════════
    s = S(VDARK)
    _title(s, "ขนาดระบบ BSIE v4.0 — ตัวเลข", WHITE)

    stats = [
        ("161", "ไฟล์โปรแกรมหลัก", GREEN), ("41", "ไฟล์หน้าจอแสดงผล", SKY),
        ("~41,000", "บรรทัดโค้ด", GOLD), ("120+", "จุดเชื่อมต่อ", BLUE),
        ("21", "ชุดคำสั่ง", MINT), ("34", "บริการย่อย", LILAC),
        ("30", "ชุดประมวลผล", ORANGE), ("20", "ตารางข้อมูล", SKY),
        ("244", "ชุดทดสอบอัตโนมัติ", GREEN), ("500+", "จุดแปลภาษา", GOLD),
        ("8", "ธนาคารไทย", RED), ("14", "ขั้นตอนประมวลผล", PURPLE),
        ("12", "รูปแบบแจ้งเตือน", ORANGE), ("5", "แบบจำลองฟอกเงิน", RED),
        ("4", "วิธีตรวจจับผิดปกติ", LILAC), ("13", "เครื่องมือสืบสวน", BLUE),
    ]
    sw_stat = (CW - GAP * 3) / 4
    sh_stat = Inches(1.15)
    for i, (num, label, clr) in enumerate(stats):
        col, row = i % 4, i // 4
        x = MX + col * (sw_stat + GAP)
        y = BODY_TOP + row * (sh_stat + Inches(0.15))
        _rect(s, x, y, sw_stat, sh_stat, DGRAY)
        _txt(s, x, y+Inches(0.08), sw_stat, Inches(0.5), num, 28, clr, True, PP_ALIGN.CENTER)
        _txt(s, x, y+Inches(0.65), sw_stat, Inches(0.35), label, 12, WHITE, False, PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════════
    # 20. ROI
    # ════════════════════════════════════════════════════════════
    s = S(GREEN)
    _title(s, "ROI — ผลตอบแทนการลงทุน", WHITE)

    roi = [
        ("เวลาวิเคราะห์", "จาก 3-5 วัน เหลือ 30 นาที ต่อคดี", 0.95, SALMON),
        ("ค่า license", "ประหยัด 742+ ล้านบาท/ปี (เทียบค่า i2 สำหรับ 1,484 สถานี)", 1.0, GOLD),
        ("ข้อผิดพลาด", "ลดข้อผิดพลาดจากวิเคราะห์ด้วยมือ — ลดการอุทธรณ์คดี", 0.75, SKY),
        ("ตรวจจับเร็ว", "ตรวจพบเร็ว ป้องกันความเสียหายล่วงหน้า", 0.85, MINT),
        ("ครอบคลุม", "1,484 สถานีตำรวจ — 77 จังหวัดทั่วประเทศ", 0.80, LILAC),
    ]
    for i, (label, desc, pct, clr) in enumerate(roi):
        y = BODY_TOP + i * Inches(1.05)
        _txt(s, MX, y, Inches(2), Inches(0.4), label, 15, WHITE, True)
        _rect(s, MX+Inches(2.2), y+Inches(0.04), CW-Inches(2.2), Inches(0.42), DGREEN)
        _rect(s, MX+Inches(2.2), y+Inches(0.04), (CW-Inches(2.2))*pct, Inches(0.42), clr)
        _txt(s, MX+Inches(2.4), y+Inches(0.04), CW-Inches(2.6), Inches(0.42), desc, 12, VDARK, True)

    _rect(s, MX, Inches(6.3), CW, Inches(0.6), DGREEN)
    _txt(s, MX, Inches(6.33), CW, Inches(0.55),
         "คืนทุนใน 1 ปี — ค่าสิทธิ์ i2 ปีเดียว มากกว่างบ SPNI ทั้ง 3 ปี", 17, WHITE, True, PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════════
    # 21. TEAM + USERS
    # ════════════════════════════════════════════════════════════
    s = S(LGRAY)
    _title(s, "ทีมพัฒนา + ประมาณการผู้ใช้งาน")

    _rect(s, MX, BODY_TOP, hw, Inches(2.3), WHITE, GREEN)
    _card_title(s, MX, BODY_TOP, hw, "เฟส 1-2: ทำได้คนเดียว + AI ช่วย", GREEN)
    _bullets(s, MX+PAD, BODY_TOP+Inches(0.5), hw-PAD*2, Inches(1.5), [
        "ผู้พัฒนาปัจจุบัน (ร.ต.อ.ณัฐวุฒิ) + **Claude AI**",
        "ผลงาน: **41,000 บรรทัด**, 244 tests, 120+ API",
        "เพียงพอสำหรับ pilot + จังหวัดนำร่อง",
    ], sz=12)

    rx = MX + hw + GAP
    _rect(s, rx, BODY_TOP, hw, Inches(2.3), WHITE, ORANGE)
    _card_title(s, rx, BODY_TOP, hw, "เฟส 3-4: ต้องการทีม 8-12 คน", ORANGE)
    _bullets(s, rx+PAD, BODY_TOP+Inches(0.5), hw-PAD*2, Inches(1.5), [
        "โปรแกรมเมอร์ฝั่งเซิร์ฟเวอร์ **2-3** │ ฝั่งหน้าจอ **1-2**",
        "ดูแลระบบ **1-2** │ วิศวกร AI **1**",
        "ทดสอบ **1-2** │ ผู้จัดการโครงการ **1**",
    ], sz=12, bc=ORANGE)

    # User estimation
    _rect(s, MX, BODY_TOP+Inches(2.6), CW, Inches(3.2), WHITE, BLUE)
    _card_title(s, MX, BODY_TOP+Inches(2.6), CW, "ประมาณการผู้ใช้งาน — จากจังหวัดสู่ทั้งประเทศ", BLUE)

    u_rows = [
        ("ระดับ", "ผู้ใช้พร้อมกัน", "ธุรกรรม/เดือน", "Storage/ปี"),
        ("Phase 1 (pilot)", "1-3 คน", "10,000", "1 GB"),
        ("Phase 2 (5 จังหวัด)", "50-150 คน", "500,000", "50 GB"),
        ("Phase 3 (20 จังหวัด)", "200-600 คน", "2,000,000", "200 GB"),
        ("Phase 4 (ทั้งประเทศ)", "2,000-5,000 คน", "10,000,000+", "1+ TB"),
    ]
    ucw = (CW - Inches(0.6)) / 4
    for i, row in enumerate(u_rows):
        y = BODY_TOP + Inches(3.2) + i * Inches(0.45)
        for j, cell in enumerate(row):
            x = MX + Inches(0.3) + j * ucw
            bg_c = NAVY if i == 0 else (WHITE if i % 2 == 1 else MGRAY)
            tc = WHITE if i == 0 else TXT
            _rect(s, x, y, ucw, Inches(0.42), bg_c)
            al = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            _txt(s, x+Inches(0.08), y+Inches(0.03), ucw-Inches(0.16), Inches(0.35),
                 cell, 11 if i > 0 else 12, tc, (i==0 or j==0), al)

    # ════════════════════════════════════════════════════════════
    # 22. CLOSING
    # ════════════════════════════════════════════════════════════
    s = S(NAVY)
    _txt(s, MX, Inches(0.4), CW, Inches(0.7), "สรุปและขั้นตอนถัดไป", 34, WHITE, True, PP_ALIGN.CENTER)
    _line(s, Inches(5.9), Inches(1.2), Inches(1.5), Pt(3), GOLD)

    closing = [
        ("1.", "BSIE v4.0 พร้อมใช้วันนี้ — ไม่ต้องรองบประมาณ", MINT),
        ("2.", "ทดลองใช้กับคดีจริง 5-10 คดี (ฟรี)", SKY),
        ("3.", "เก็บความเห็นจากพนักงานสอบสวน", GOLD),
        ("4.", "เสนองบเฟส 2: ~8 ล้านบาท (5 จังหวัดนำร่อง)", SALMON),
        ("5.", "ขยายสู่ระบบ SPNI ทั้งประเทศ — 77 จังหวัด 1,484 สถานี", LILAC),
    ]
    for i, (num, text, clr) in enumerate(closing):
        y = Inches(1.6) + i * Inches(0.8)
        _txt(s, Inches(3.2), y, Inches(0.6), Inches(0.55), num, 24, clr, True)
        _txt(s, Inches(3.8), y + Inches(0.05), Inches(7), Inches(0.5), text, 18, WHITE)

    _rect(s, MX, Inches(5.7), CW, Inches(0.7), CDARK)
    _txt(s, MX, Inches(5.73), CW, Inches(0.65),
         "งบ SPNI ทั้ง 3 ปี (~170 ล้าน) น้อยกว่า ค่าสิทธิ์ i2 ปีเดียว (742+ ล้าน) คืนทุนทันที",
         15, GOLD, True, PP_ALIGN.CENTER)

    _txt(s, MX, Inches(6.7), CW, Inches(0.35),
         "ร้อยตำรวจเอกณัฐวุฒิ สาหร่ายทอง │ โทร: 096-776-8757 │ BSIE v4.0 — Bank Statement Intelligence Engine",
         11, SUBTLE, False, PP_ALIGN.CENTER)

    # ── Save ──
    out = Path(__file__).parent.parent / "docs" / "SPNI_BSIE_Presentation.pptx"
    prs.save(str(out))
    print(f"✅ Saved: {out}")
    print(f"   {len(prs.slides)} slides")
    return out


if __name__ == "__main__":
    generate()
